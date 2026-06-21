from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_FILE = ROOT / "FraudX_Viva_Presentation.pptx"


def _set_title(slide, title: str) -> None:
    title_shape = slide.shapes.title
    title_shape.text = title
    p = title_shape.text_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)


def _set_bullets(slide, bullets: list[str], font_size: int = 22) -> None:
    # Use "Title and Content" layout; content placeholder index 1
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(40, 40, 40)


def _add_slide(prs: Presentation, title: str, bullets: list[str], font_size: int = 22) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    _set_title(slide, title)
    _set_bullets(slide, bullets, font_size=font_size)


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = "FraudX — Advanced Fraud Detection System"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Viva Presentation • {date.today().strftime('%d %b %Y')}\nPython • Streamlit • ML (RandomForest) • Explainable Outputs"

    # Light touch formatting
    t = slide.shapes.title.text_frame.paragraphs[0].font
    t.size = Pt(40)
    t.bold = True

    s = subtitle.text_frame.paragraphs[0].font
    s.size = Pt(18)


def build_ppt(out_path: Path = OUT_FILE) -> Path:
    prs = Presentation()
    prs.core_properties.title = "FraudX Viva Presentation"
    prs.core_properties.subject = "Architecture, Algorithms, Interfaces, Roadmap"

    _add_title_slide(prs)

    _add_slide(
        prs,
        "Problem & Objective",
        [
            "Goal: detect suspicious financial transactions with explainable reasons",
            "Support real-time-like scoring on uploaded/batch CSV data",
            "Provide a simple UI for threshold tuning and analysis",
        ],
    )

    _add_slide(
        prs,
        "System Architecture (High Level)",
        [
            "Presentation: Streamlit app (`app.py`) for upload, threshold, visualizations, export",
            "Detection engine: preprocessing → model scoring → thresholding → fraud-type rules",
            "Artifacts: trained model + scaler + feature importances (`*.joblib`)",
            "Optional interface: FastAPI service (`api/main.py`) for programmatic scoring",
        ],
        font_size=20,
    )

    _add_slide(
        prs,
        "End-to-End Flow",
        [
            "Input: upload CSV (or generate synthetic data via `main.py`)",
            "Preprocess: align columns to training schema, encode + scale features",
            "Score: RandomForest `predict_proba` → `Suspicion_Score`",
            "Decide: user-selected threshold → `Is_Suspicious`",
            "Explain: generate `Suspicion_Reasons` + assign `Fraud_Type` (rules)",
            "Output: dashboard + downloadable `suspicious_transactions.csv`",
        ],
        font_size=19,
    )

    _add_slide(
        prs,
        "Algorithm Development — ML Scoring",
        [
            "Model: RandomForest classifier (trained on synthetic dataset)",
            "Output: probability of fraud used as `Suspicion_Score` (0.0–1.0)",
            "Why RF: handles non-linear patterns, robust baseline, easy inference",
            "Artifacts saved for reproducible inference (model + scaler + feature map)",
        ],
        font_size=20,
    )

    _add_slide(
        prs,
        "Algorithm Development — Explainability & Labels",
        [
            "`explain_suspicion`: human-readable reasons using important features + history",
            "`classify_fraud_type`: rule-based primary label for suspicious cases",
            "Examples: Known Fraudster, Sanctioned Entity, ATO Risk, Rapid Large Transfers",
            "Design choice: explanations are part of the output (not a separate report)",
        ],
        font_size=20,
    )

    _add_slide(
        prs,
        "Interfaces Identified",
        [
            "UI interface (Streamlit): upload CSV, set threshold slider, view charts, export results",
            "Core internal interfaces:",
            "• `preprocess_data(df) -> X_aligned_scaled`",
            "• `detect_suspicious_transactions(df, threshold) -> df_with_scores_labels`",
            "• `explain_suspicion(row/context) -> reasons`",
            "Optional API interface (FastAPI): health + single/batch detection endpoints",
        ],
        font_size=18,
    )

    _add_slide(
        prs,
        "Performance & Evaluation (Current)",
        [
            "Dataset: 10,000+ synthetic transactions, 20+ engineered features",
            "Reported test performance (current):",
            "• Accuracy: 95.2%  • Precision: 94.8%  • Recall: 93.5%  • F1: 94.1%",
            "Runtime: designed for fast scoring on uploaded/batch data",
        ],
        font_size=20,
    )

    _add_slide(
        prs,
        "Components Implemented So Far",
        [
            "Synthetic data generation + model training (`main.py`)",
            "Saved artifacts (`fraud_detection_model.joblib`, scaler, feature importances)",
            "Streamlit dashboard: KPI cards, charts, details panel, CSV export (`app.py`)",
            "Detection + explanation module (`suspicious_by_model.py`)",
            "Optional FastAPI scoring endpoints (`api/main.py`)",
        ],
        font_size=20,
    )

    _add_slide(
        prs,
        "Planned Enhancements (Next Steps)",
        [
            "Add richer real-world data support + schema validation and tests",
            "Improve threshold selection using PR curve / cost-based tuning",
            "Model monitoring: drift checks and periodic retraining",
            "Better explainability: SHAP for per-transaction feature contributions",
            "Deployment: containerize + CI checks (optional)",
        ],
        font_size=20,
    )

    _add_slide(
        prs,
        "Viva Q&A — Ready Answers",
        [
            "Why this architecture? Modular UI/engine/artifacts enables easy iteration",
            "Why RandomForest? Strong baseline for non-linear patterns; fast inference",
            "How is fraud decided? `Suspicion_Score` from model + user threshold slider",
            "How do you explain flags? `Suspicion_Reasons` + rule-based `Fraud_Type`",
            "How will you scale? API service + stateless scoring + batch/stream integration",
        ],
        font_size=19,
    )

    # Small footer on each slide (except title) for a clean simple look
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.4))
        tf = tx.text_frame
        tf.text = "FraudX • Architecture • Algorithms • Interfaces • Roadmap"
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(120, 120, 120)

    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    out = build_ppt()
    print(f"Created: {out}")
